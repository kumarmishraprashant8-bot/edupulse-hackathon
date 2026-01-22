"""PPTX generation service for micro-modules and LFA exports."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from typing import List, Dict


class PPTXGenerator:
    """Generate PowerPoint presentations for modules and LFA."""
    
    def __init__(self, exports_path: str = "exports"):
        """Initialize with exports directory."""
        self.exports_path = exports_path
        os.makedirs(exports_path, exist_ok=True)
    
    def generate_micro_module(
        self,
        title: str,
        topic: str,
        advice: str,
        materials: str,
        cluster: str,
        output_filename: str
    ) -> str:
        """
        Generate 2-slide micro-module PPTX.
        
        Args:
            title: Module title
            topic: Topic tag
            advice: Advice text (will be split into bullets)
            materials: Required materials
            cluster: Target cluster
            output_filename: Output file name (without path)
        
        Returns:
            Full path to generated file
        """
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: Title + Action Steps
        slide1 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        
        # Title
        title_shape = slide1.shapes.title
        title_shape.text = title
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        
        # Subtitle with cluster
        subtitle = slide1.placeholders[1]
        subtitle.text = f"For: {cluster}\nTopic: {topic.replace('-', ' ').title()}"
        
        # Content - action steps
        content = slide1.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(8), Inches(4)
        )
        text_frame = content.text_frame
        text_frame.word_wrap = True
        
        # Add heading
        p = text_frame.paragraphs[0]
        p.text = "Quick Action Steps:"
        p.font.size = Pt(24)
        p.font.bold = True
        p.space_after = Pt(12)
        
        # Split advice into steps
        steps = [s.strip() for s in advice.split('\n') if s.strip()]
        for step in steps:
            p = text_frame.add_paragraph()
            p.text = step
            p.font.size = Pt(18)
            p.space_after = Pt(8)
            p.level = 0
        
        # Slide 2: Classroom Script + Resources
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Title
        title2 = slide2.shapes.title
        title2.text = "Implementation Guide"
        
        # Left box - script
        left_box = slide2.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(4.5), Inches(5)
        )
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        
        p = left_frame.paragraphs[0]
        p.text = "Sample Classroom Script:"
        p.font.size = Pt(20)
        p.font.bold = True
        p.space_after = Pt(12)
        
        script_text = self._generate_script(topic, steps[0] if steps else advice)
        p = left_frame.add_paragraph()
        p.text = script_text
        p.font.size = Pt(16)
        
        # Right box - resources
        right_box = slide2.shapes.add_textbox(
            Inches(5.5), Inches(1.5), Inches(4), Inches(5)
        )
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        
        p = right_frame.paragraphs[0]
        p.text = "Materials Needed:"
        p.font.size = Pt(20)
        p.font.bold = True
        p.space_after = Pt(12)
        
        p = right_frame.add_paragraph()
        p.text = f"• {materials}"
        p.font.size = Pt(16)
        p.space_after = Pt(12)
        
        p = right_frame.add_paragraph()
        p.text = "Time Required:"
        p.font.size = Pt(18)
        p.font.bold = True
        p.space_after = Pt(8)
        
        p = right_frame.add_paragraph()
        p.text = "• 15-20 minutes"
        p.font.size = Pt(16)
        p.space_after = Pt(12)
        
        p = right_frame.add_paragraph()
        p.text = "Support Available:"
        p.font.size = Pt(18)
        p.font.bold = True
        p.space_after = Pt(8)
        
        p = right_frame.add_paragraph()
        p.text = "• Contact your CRP\n• WhatsApp support\n• Demo video link"
        p.font.size = Pt(16)
        
        # Save
        output_path = os.path.join(self.exports_path, output_filename)
        prs.save(output_path)
        return output_path
    
    def generate_lfa_export(
        self,
        title: str,
        problem_statement: str,
        student_change: str,
        stakeholders: List[str],
        practice_changes: List[str],
        indicators: List[str],
        output_filename: str
    ) -> str:
        """
        Generate 1-2 slide LFA framework export.
        
        Returns:
            Full path to generated file
        """
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: Overview
        slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
        
        # Title
        title_box = slide1.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Create grid layout
        # Problem (top left)
        prob_box = slide1.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(4.5), Inches(2.5)
        )
        self._add_lfa_section(
            prob_box.text_frame,
            "Problem",
            problem_statement
        )
        
        # Desired Change (top right)
        change_box = slide1.shapes.add_textbox(
            Inches(5.5), Inches(1.5), Inches(4), Inches(2.5)
        )
        self._add_lfa_section(
            change_box.text_frame,
            "Desired Student Change",
            student_change
        )
        
        # Stakeholders (bottom left)
        stake_box = slide1.shapes.add_textbox(
            Inches(0.5), Inches(4.5), Inches(4.5), Inches(2.5)
        )
        self._add_lfa_section(
            stake_box.text_frame,
            "Key Stakeholders",
            "\n".join([f"• {s}" for s in stakeholders])
        )
        
        # Practice Changes (bottom right)
        practice_box = slide1.shapes.add_textbox(
            Inches(5.5), Inches(4.5), Inches(4), Inches(2.5)
        )
        self._add_lfa_section(
            practice_box.text_frame,
            "Practice Changes",
            "\n".join([f"• {p}" for p in practice_changes[:3]])
        )
        
        # Slide 2: Indicators
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        
        title2 = slide2.shapes.title
        title2.text = "Key Indicators & Measurement"
        
        indicator_box = slide2.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(4.5)
        )
        ind_frame = indicator_box.text_frame
        ind_frame.word_wrap = True
        
        p = ind_frame.paragraphs[0]
        p.text = "Success Indicators:"
        p.font.size = Pt(24)
        p.font.bold = True
        p.space_after = Pt(12)
        
        for i, indicator in enumerate(indicators, 1):
            p = ind_frame.add_paragraph()
            p.text = f"{i}. {indicator}"
            p.font.size = Pt(18)
            p.space_after = Pt(10)
        
        # Save
        output_path = os.path.join(self.exports_path, output_filename)
        prs.save(output_path)
        return output_path
    
    def _add_lfa_section(self, text_frame, heading: str, content: str):
        """Add formatted section to LFA slide."""
        p = text_frame.paragraphs[0]
        p.text = heading
        p.font.size = Pt(18)
        p.font.bold = True
        p.space_after = Pt(8)
        
        p = text_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(14)
    
    def _generate_script(self, topic: str, first_step: str) -> str:
        """Generate sample classroom script based on topic."""
        scripts = {
            "subtraction-borrowing": (
                '"Today we\'ll practice subtraction with borrowing. '
                'Take out your pebbles. Let\'s show 13-7 together. '
                'Count 13 pebbles. Now, can we take away 7 from the 3 we have? '
                'No! So we need to borrow from the tens place..."'
            ),
            "fractions-conceptual": (
                '"Let\'s explore fractions! Take your paper and fold it in half. '
                'How many equal parts? That\'s right - 2 parts. Each part is 1/2. '
                'Now fold again. How many parts now? 4 parts - each is 1/4..."'
            ),
            "multiplication-tables": (
                '"Let\'s sing the 2s! 2, 4, 6, 8... Now let\'s show it with dots. '
                'Draw 2 rows of 4 dots. How many total? Count with me: 2, 4, 6, 8!"'
            ),
        }
        return scripts.get(topic, f'"Let\'s start: {first_step}"')